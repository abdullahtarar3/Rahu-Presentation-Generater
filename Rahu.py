import json
import os
from pptx import Presentation
import tkinter as tk
from tkinter import messagebox, filedialog, ttk
import requests  # For Serper API

# Function to set API key based on user selection
def set_api_key(api_choice, api_key_path):
    if not api_key_path:
        raise ValueError("No API key file selected")
    
    # Load the API key from the selected .env file
    with open(api_key_path, 'r') as file:
        for line in file:
            if api_choice.upper() in line:
                return line.split('=')[1].strip()
    raise ValueError(f"API key for {api_choice} not found in the selected file")

# Function to fetch fresh data using Serper API
def fetch_fresh_data(query):
    serper_api_key = set_api_key("Serper", api_path_entry.get())  # Fetch Serper API key
    url = "https://google.serper.dev/search"
    payload = json.dumps({"q": query})
    headers = {
        'X-API-KEY': serper_api_key,
        'Content-Type': 'application/json'
    }
    response = requests.post(url, headers=headers, data=payload)
    if response.status_code == 200:
        return response.json()
    else:
        raise ValueError(f"Failed to fetch fresh data: {response.text}")

def generate_presentation():
    # Collect user inputs
    presentation_title = title_entry.get()
    specific_topic = topic_entry.get()
    language = language_entry.get()
    role = role_entry.get()
    writing_tone = tone_entry.get()
    context = context_entry.get()
    api_choice = api_var.get()
    api_key_path = api_path_entry.get()
    output_directory = output_path_entry.get()
    num_slides = slides_entry.get()
    use_fresh_data = fresh_data_var.get()  # Check if user wants fresh data

    if not presentation_title:
        messagebox.showwarning("Input Error", "Please enter a presentation title!")
        return

    if not num_slides.isdigit() or int(num_slides) <= 0:
        messagebox.showwarning("Input Error", "Please enter a valid number of slides!")
        return

    # Ask user for API key
    try:
        api_key = set_api_key(api_choice, api_key_path)
    except ValueError as e:
        messagebox.showerror("Error", str(e))
        return

    if not output_directory:
        messagebox.showwarning("Input Error", "Please select an output directory!")
        return

    # Build the prompt dynamically
    question = (
        f"Generate a {num_slides}-slide presentation on '{presentation_title}'. "
        "Each slide should have a header and content (50-60 words per slide). "
        "The final slide should contain discussion questions. "
    )

    if specific_topic:
        question += f" Focus specifically on '{specific_topic}'."
    if language:
        question += f" Write in {language}."
    if role:
        question += f" Assume the role of a {role} while writing."
    if writing_tone:
        question += f" Use a {writing_tone} tone."
    if context:
        question += f" Context: {context}"

    # If user wants fresh data, fetch it using Serper API
    if use_fresh_data:
        try:
            fresh_data = fetch_fresh_data(presentation_title)
            question += f" Use the following fresh data: {json.dumps(fresh_data)}"
        except Exception as e:
            messagebox.showerror("Error", f"Failed to fetch fresh data: {str(e)}")
            return

    query_json = '''{
        "input_text": "[[QUERY]]",
        "output_format": "json",
        "json_structure": {
            "slides": [
                {
                    "header": "{{header}}",
                    "content": "{{content}}"
                }
            ]
        }
    }'''

    # Replace query placeholder
    prompt = query_json.replace("[[QUERY]]", question)

    try:
        # Generate response from the selected API
        if api_choice == "OpenAI":
            import openai
            openai.api_key = api_key
            completion = openai.ChatCompletion.create(
                model="gpt-3.5-turbo",
                messages=[{"role": "user", "content": prompt}]
            )
            response = completion.choices[0].message.content
        elif api_choice == "Gemini":
            # Corrected import and method for Gemini API
            import google.generativeai as genai
            genai.configure(api_key=api_key)
            model = genai.GenerativeModel(model_name='gemini-1.5-flash')
            response = model.generate_content(prompt).text

            # Check if the response is empty
            if not response:
                raise ValueError("Received empty response from Gemini API")
            
        elif api_choice == "DeepSeek":
            from deepseek import DeepSeek
            client = DeepSeek(api_key)
            response = client.completions.create(
                model="deepseek-chat",
                prompt=question,
                max_tokens=128
            ).choices[0].text
        elif api_choice == "Kimi":
            from openai import Client
            client = Client(
                api_key=api_key,
                base_url="https://api.moonshot.cn/v1",
            )
            completion = client.chat.completions.create(
                model="moonshot-v1-8k",
                messages=[{"role": "user", "content": prompt}]
            )
            response = completion.choices[0].message.content
        else:
            raise ValueError("Unsupported API choice")

        # Parse JSON response
        if api_choice == "Gemini":
            # Gemini API returns a string, so we need to parse it as JSON
            try:
                r = json.loads(response)
            except json.JSONDecodeError:
                raise ValueError("Failed to parse JSON response from Gemini API")
        else:
            r = json.loads(response)
        slide_data = r["slides"]
    except json.JSONDecodeError as e:
        messagebox.showerror("Error", f"Failed to parse JSON response: {str(e)}")
        return
    except Exception as e:
        messagebox.showerror("Error", f"Failed to generate presentation: {str(e)}")
        return

    # Create PowerPoint Presentation
    prs = Presentation()

    for slide in slide_data:
        slide_layout = prs.slide_layouts[1]  # Title & Content Layout
        new_slide = prs.slides.add_slide(slide_layout)

        # Set slide title
        if "header" in slide and slide["header"]:
            title = new_slide.shapes.title
            title.text = slide["header"]

        # Set slide content
        if "content" in slide and slide["content"]:
            text_box = new_slide.shapes.placeholders[1]
            text_frame = text_box.text_frame
            text_frame.text = slide["content"]

    # Save the presentation in the selected directory with the presentation title as the filename
    output_path = os.path.join(output_directory, f"{presentation_title}.pptx")
    prs.save(output_path)

    messagebox.showinfo("Success", f"Presentation saved at: {output_path}")

# GUI Setup
root = tk.Tk()
root.title("Rahu Presentation Generator")
root.geometry("800x800")
root.configure(bg="#f0f0f0")

# Style configuration
style = {
    "bg": "#f0f0f0",
    "font": ("Arial", 12),
    "padx": 10,
    "pady": 5
}

entry_style = {
    "width": 50,
    "font": ("Arial", 10)
}

button_style = {
    "font": ("Arial", 10),
    "bg": "#4CAF50",
    "fg": "white",
    "padx": 10,
    "pady": 5
}

# Main Frame
main_frame = ttk.Frame(root, padding="10")
main_frame.grid(row=0, column=0, sticky=(tk.W, tk.E, tk.N, tk.S))

# Title label and entry
ttk.Label(main_frame, text="Enter Presentation Title").grid(row=0, column=0, sticky=tk.W, pady=5)
title_entry = ttk.Entry(main_frame, width=50)
title_entry.grid(row=0, column=1, pady=5)

# Topic label and entry
ttk.Label(main_frame, text="Specific Topic or Sub-topic (Optional)").grid(row=1, column=0, sticky=tk.W, pady=5)
topic_entry = ttk.Entry(main_frame, width=50)
topic_entry.grid(row=1, column=1, pady=5)

# Language label and entry
ttk.Label(main_frame, text="Language Option (Optional)").grid(row=2, column=0, sticky=tk.W, pady=5)
language_entry = ttk.Entry(main_frame, width=50)
language_entry.grid(row=2, column=1, pady=5)

# Role label and entry
ttk.Label(main_frame, text="Role (e.g., Professor, Student, etc.) (Optional)").grid(row=3, column=0, sticky=tk.W, pady=5)
role_entry = ttk.Entry(main_frame, width=50)
role_entry.grid(row=3, column=1, pady=5)

# Writing Tone label and entry
ttk.Label(main_frame, text="Writing Tone (e.g., Formal, Casual, Persuasive) (Optional)").grid(row=4, column=0, sticky=tk.W, pady=5)
tone_entry = ttk.Entry(main_frame, width=50)
tone_entry.grid(row=4, column=1, pady=5)

# Context label and entry
ttk.Label(main_frame, text="Context (Optional)").grid(row=5, column=0, sticky=tk.W, pady=5)
context_entry = ttk.Entry(main_frame, width=50)
context_entry.grid(row=5, column=1, pady=5)

# Number of Slides label and entry
ttk.Label(main_frame, text="Number of Slides").grid(row=6, column=0, sticky=tk.W, pady=5)
slides_entry = ttk.Entry(main_frame, width=50)
slides_entry.grid(row=6, column=1, pady=5)

# Fresh Data Checkbox
fresh_data_var = tk.BooleanVar()
fresh_data_checkbox = ttk.Checkbutton(main_frame, text="Use Fresh Data (Optional)", variable=fresh_data_var)
fresh_data_checkbox.grid(row=7, column=0, sticky=tk.W, pady=5)

# API choice label and dropdown
ttk.Label(main_frame, text="Select API").grid(row=8, column=0, sticky=tk.W, pady=5)
api_var = tk.StringVar(value="OpenAI")
api_dropdown = ttk.Combobox(main_frame, textvariable=api_var, values=["OpenAI", "Gemini", "DeepSeek", "Kimi"], width=47)
api_dropdown.grid(row=8, column=1, pady=5)

# API Path label and entry
ttk.Label(main_frame, text="Path to API Key (.env file)").grid(row=9, column=0, sticky=tk.W, pady=5)
api_path_entry = ttk.Entry(main_frame, width=50)
api_path_entry.grid(row=9, column=1, pady=5)

# API Path browse button
def browse_api_path():
    api_path = filedialog.askopenfilename(title="Select .env file for API")
    api_path_entry.delete(0, tk.END)
    api_path_entry.insert(0, api_path)

browse_api_button = ttk.Button(main_frame, text="Browse", command=browse_api_path)
browse_api_button.grid(row=9, column=2, pady=5)

# Output Path label and entry
ttk.Label(main_frame, text="Output Directory").grid(row=10, column=0, sticky=tk.W, pady=5)
output_path_entry = ttk.Entry(main_frame, width=50)
output_path_entry.grid(row=10, column=1, pady=5)

# Output Path browse button
def browse_output_path():
    output_path = filedialog.askdirectory(title="Select Output Directory")
    output_path_entry.delete(0, tk.END)
    output_path_entry.insert(0, output_path)

browse_output_button = ttk.Button(main_frame, text="Browse", command=browse_output_path)
browse_output_button.grid(row=10, column=2, pady=5)

# Generate button
generate_button = ttk.Button(main_frame, text="Generate Presentation", command=generate_presentation)
generate_button.grid(row=11, column=1, pady=20)
import webbrowser
# Function to open Instagram profile
def open_instagram(event):
    webbrowser.open("https://linktr.ee/abdullahtarar.3")

# Function to open developer info page
def open_dev_info(event):
    webbrowser.open("https://linktr.ee/abdullahtarar.3")

# Developer information
developer_info = ttk.Label(main_frame, text="Developed by Abdullah Tarar", font=("Arial", 10), foreground="blue", cursor="hand2")
developer_info.grid(row=12, column=1, pady=10)
developer_info.bind("<Button-1>", open_instagram)

# Additional developer info
more_info = ttk.Label(main_frame, text="Get more info about the developer", font=("Arial", 10), foreground="blue", cursor="hand2")
more_info.grid(row=13, column=1, pady=5)
more_info.bind("<Button-1>", open_dev_info)



# Run the GUI loop
root.mainloop()
